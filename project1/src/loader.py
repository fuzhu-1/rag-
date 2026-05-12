"""
Enterprise-RAG: Multi-format document loader with OCR and multimodal support.
"""
import hashlib
import os
from pathlib import Path
from typing import Any

import fitz  # PyMuPDF
from docx import Document as DocxDocument
from loguru import logger
from PIL import Image

from src.config import config


class Document:
    """Represents a loaded document with metadata."""

    def __init__(
        self,
        content: str,
        metadata: dict[str, Any] | None = None,
    ):
        self.content = content
        self.metadata = metadata or {}
        self.doc_id = hashlib.md5(
            f"{metadata.get('source', '')}{content[:200]}".encode()
        ).hexdigest()[:16]


class DocumentLoader:
    """Load and parse documents from various formats."""

    ALLOWED_EXTENSIONS: set[str] = {
        ".pdf", ".docx", ".doc", ".txt", ".md",
        ".png", ".jpg", ".jpeg", ".xlsx", ".xls", ".pptx", ".csv",
    }

    def __init__(self, data_dir: str | None = None):
        cfg = config
        self.data_dir = Path(data_dir or cfg["project"]["data_dir"])
        self.allowed_exts = set(
            cfg.get("security", {}).get("allowed_extensions", self.ALLOWED_EXTENSIONS)
        )
        self.max_file_size = (
            cfg.get("security", {}).get("max_file_size_mb", 100) * 1024 * 1024
        )
        self.ocr_lang = cfg.get("multimodal", {}).get("ocr_lang", "chi_sim+eng")
        self.multimodal_enabled = cfg.get("multimodal", {}).get("enabled", False)

    def load_directory(self, directory: str | None = None) -> list[Document]:
        """Load all supported documents from a directory."""
        target = Path(directory) if directory else self.data_dir
        if not target.exists():
            logger.warning(f"Directory not found: {target}")
            return []

        documents: list[Document] = []
        for file_path in target.iterdir():
            if file_path.is_file() and file_path.suffix.lower() in self.allowed_exts:
                docs = self.load_file(str(file_path))
                documents.extend(docs)
                logger.info(f"Loaded: {file_path.name} ({len(docs)} segments)")

        return documents

    def load_file(self, file_path: str) -> list[Document]:
        """Load a single file, dispatching by extension."""
        path = Path(file_path)
        self._validate_file(path)

        suffix = path.suffix.lower()
        if suffix == ".pdf":
            return self._load_pdf(path)
        elif suffix in (".docx", ".doc"):
            return self._load_docx(path)
        elif suffix in (".txt", ".md"):
            return self._load_text(path)
        elif suffix in (".png", ".jpg", ".jpeg"):
            return self._load_image(path)
        elif suffix in (".xlsx", ".xls"):
            return self._load_excel(path)
        elif suffix == ".csv":
            return self._load_csv(path)
        elif suffix == ".pptx":
            return self._load_pptx(path)
        else:
            raise ValueError(f"Unsupported file type: {suffix}")

    def _validate_file(self, path: Path) -> None:
        """Security: validate file extension and prevent path traversal."""
        resolved = path.resolve()
        if ".." in str(path):
            raise ValueError("Path traversal detected")

        if not resolved.exists():
            raise FileNotFoundError(f"File not found: {path}")

        if path.suffix.lower() not in self.allowed_exts:
            raise ValueError(f"File extension not allowed: {path.suffix}")

        if resolved.stat().st_size > self.max_file_size:
            raise ValueError(
                f"File too large: {resolved.stat().st_size / 1024 / 1024:.1f}MB "
                f"(max {self.max_file_size / 1024 / 1024:.0f}MB)"
            )

    def _load_pdf(self, path: Path) -> list[Document]:
        """Load PDF with PyMuPDF, including OCR for image-based pages."""
        documents: list[Document] = []
        doc = fitz.open(str(path))

        for page_num, page in enumerate(doc, start=1):
            text = page.get_text()

            # If page has no extractable text, try OCR
            if not text.strip():
                text = self._ocr_page(page)

            # Extract images from page
            images_text = ""
            if self.multimodal_enabled:
                images_text = self._extract_page_images(page, page_num)

            full_text = (text + "\n" + images_text).strip()
            if full_text:
                documents.append(Document(
                    content=full_text,
                    metadata={
                        "source": path.name,
                        "file_path": str(path.resolve()),
                        "page": page_num,
                        "total_pages": doc.page_count,
                        "type": "pdf",
                        "has_images": bool(images_text),
                    },
                ))

        doc.close()
        return documents

    def _ocr_page(self, page: fitz.Page) -> str:
        """OCR a PDF page using pytesseract."""
        try:
            import pytesseract

            pix = page.get_pixmap(dpi=200)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            text = pytesseract.image_to_string(img, lang=self.ocr_lang)
            return text.strip()
        except ImportError:
            logger.warning("pytesseract not installed, skipping OCR")
            return ""
        except Exception as e:
            logger.error(f"OCR failed: {e}")
            return ""

    def _extract_page_images(self, page: fitz.Page, page_num: int) -> str:
        """Extract images from a PDF page and generate descriptions."""
        descriptions: list[str] = []
        image_list = page.get_images(full=True)

        for img_idx, img_info in enumerate(image_list):
            try:
                xref = img_info[0]
                base_image = page.parent.extract_image(xref)
                img_bytes = base_image["image"]
                img_ext = base_image["ext"]
                img_path = Path(self.data_dir) / f"_images/page{page_num}_img{img_idx}.{img_ext}"
                img_path.parent.mkdir(parents=True, exist_ok=True)
                img_path.write_bytes(img_bytes)

                # If multimodal enabled, generate image description
                if self.multimodal_enabled:
                    desc = self._describe_image(str(img_path))
                    if desc:
                        descriptions.append(f"[图片内容 page{page_num} img{img_idx}]: {desc}")
            except Exception as e:
                logger.warning(f"Image extraction failed on page {page_num}: {e}")

        return "\n".join(descriptions)

    def _describe_image(self, image_path: str) -> str:
        """Generate image description using vision model."""
        try:
            from openai import OpenAI

            vision_cfg = config.get("multimodal", {})
            model = vision_cfg.get("vision_model", "gpt-4o-mini")

            if model.startswith("gpt"):
                client = OpenAI(
                    api_key=os.environ.get("OPENAI_API_KEY", ""),
                    base_url=config.get("llm", {}).get("openai", {}).get(
                        "api_base", "https://api.openai.com/v1"
                    ),
                )
                import base64

                with open(image_path, "rb") as f:
                    img_b64 = base64.b64encode(f.read()).decode()

                response = client.chat.completions.create(
                    model=model,
                    messages=[{
                        "role": "user",
                        "content": [
                            {
                                "type": "text",
                                "text": "请详细描述这张图片的内容，包括图表、数据、文字等所有信息。",
                            },
                            {
                                "type": "image_url",
                                "image_url": {"url": f"data:image/png;base64,{img_b64}"},
                            },
                        ],
                    }],
                    max_tokens=500,
                )
                return response.choices[0].message.content or ""
            return ""
        except Exception as e:
            logger.warning(f"Image description failed: {e}")
            return ""

    def _load_docx(self, path: Path) -> list[Document]:
        """Load Word documents."""
        doc = DocxDocument(str(path))
        full_text: list[str] = []

        for para in doc.paragraphs:
            if para.text.strip():
                full_text.append(para.text)

        # Extract tables
        for table_idx, table in enumerate(doc.tables):
            table_text: list[str] = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                table_text.append(" | ".join(cells))
            full_text.append(f"\n[表格 {table_idx + 1}]:\n" + "\n".join(table_text))

        content = "\n".join(full_text)
        if content.strip():
            return [Document(
                content=content,
                metadata={
                    "source": path.name,
                    "file_path": str(path.resolve()),
                    "type": "docx",
                    "paragraphs": len(doc.paragraphs),
                },
            )]
        return []

    def _load_text(self, path: Path) -> list[Document]:
        """Load plain text and Markdown files."""
        content = path.read_text(encoding="utf-8", errors="replace")
        if content.strip():
            return [Document(
                content=content,
                metadata={
                    "source": path.name,
                    "file_path": str(path.resolve()),
                    "type": path.suffix.lstrip("."),
                },
            )]
        return []

    def _load_image(self, path: Path) -> list[Document]:
        """Load standalone image with OCR + description."""
        documents: list[Document] = []

        # OCR
        try:
            import pytesseract
            img = Image.open(str(path))
            ocr_text = pytesseract.image_to_string(img, lang=self.ocr_lang)
        except ImportError:
            ocr_text = ""
        except Exception as e:
            logger.error(f"OCR failed for {path.name}: {e}")
            ocr_text = ""

        # Vision description
        desc = ""
        if self.multimodal_enabled:
            desc = self._describe_image(str(path))

        content_parts = []
        if ocr_text.strip():
            content_parts.append(f"[OCR文本]: {ocr_text.strip()}")
        if desc:
            content_parts.append(f"[图片描述]: {desc}")

        if content_parts:
            documents.append(Document(
                content="\n".join(content_parts),
                metadata={
                    "source": path.name,
                    "file_path": str(path.resolve()),
                    "type": "image",
                    "has_ocr": bool(ocr_text.strip()),
                    "has_description": bool(desc),
                },
            ))

        return documents

    def _load_excel(self, path: Path) -> list[Document]:
        """Load Excel spreadsheets."""
        import openpyxl

        wb = openpyxl.load_workbook(str(path), data_only=True)
        all_text: list[str] = [f"# 工作表列表: {wb.sheetnames}"]

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows_text: list[str] = [f"\n## 工作表: {sheet_name}"]
            for row in ws.iter_rows(values_only=True):
                row_str = " | ".join(str(cell) if cell is not None else "" for cell in row)
                if row_str.strip():
                    rows_text.append(row_str)
            all_text.extend(rows_text)

        content = "\n".join(all_text)
        if content.strip():
            return [Document(
                content=content,
                metadata={
                    "source": path.name,
                    "file_path": str(path.resolve()),
                    "type": "excel",
                    "sheets": wb.sheetnames,
                },
            )]
        return []

    def _load_csv(self, path: Path) -> list[Document]:
        """Load CSV files."""
        import csv

        rows: list[str] = []
        with open(path, "r", encoding="utf-8-sig", errors="replace") as f:
            reader = csv.reader(f)
            header = next(reader, None)
            if header:
                rows.append(" | ".join(header))
                rows.append("-" * len(" | ".join(header)))
            for row in reader:
                rows.append(" | ".join(row))

        content = "\n".join(rows)
        if content.strip():
            return [Document(
                content=content,
                metadata={
                    "source": path.name,
                    "file_path": str(path.resolve()),
                    "type": "csv",
                },
            )]
        return []

    def _load_pptx(self, path: Path) -> list[Document]:
        """Load PowerPoint presentations."""
        from pptx import Presentation

        prs = Presentation(str(path))
        slides_text: list[str] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_content: list[str] = [f"\n# 幻灯片 {slide_num}"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_content.append(para.text)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        slide_content.append(" | ".join(cells))
            slides_text.append("\n".join(slide_content))

        content = "\n".join(slides_text)
        if content.strip():
            return [Document(
                content=content,
                metadata={
                    "source": path.name,
                    "file_path": str(path.resolve()),
                    "type": "pptx",
                    "slides": len(prs.slides),
                },
            )]
        return []
